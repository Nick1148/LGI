#!/usr/bin/env python3
"""MCP 서버 허브 사업 설명 PPT 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === Color Palette ===
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1A, 0x25, 0x3C)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, default_size=16, default_color=TEXT_LIGHT):
    """lines: list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, size, color, bold, align = line_data, default_size, default_color, False, PP_ALIGN.LEFT
        else:
            text = line_data[0]
            size = line_data[1] if len(line_data) > 1 else default_size
            color = line_data[2] if len(line_data) > 2 else default_color
            bold = line_data[3] if len(line_data) > 3 else False
            align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        p.space_after = Pt(4)
    return txBox


# ================================================================
# SLIDE 1: 표지
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
         "한국 로컬 MCP 서버 허브", 48, TEXT_WHITE, True, PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(3), Inches(11), Inches(0.8),
         "AI 에이전트가 한국 서비스를 사용할 수 있게 연결하는 인프라 플랫폼", 24, TEXT_GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
         "제로클릭 시대, AI가 사용하는 서비스를 만들자", 20, ACCENT_BLUE, True, PP_ALIGN.CENTER)

# 하단 키워드 카드들
keywords = [
    ("MCP Protocol", ACCENT_BLUE),
    ("Korean Local API", ACCENT_PURPLE),
    ("AI Agent Infra", ACCENT_GREEN),
    ("1인 개발자", ACCENT_ORANGE),
]
card_w = Inches(2.5)
start_x = Inches(1.5)
for i, (kw, color) in enumerate(keywords):
    x = start_x + i * Inches(2.8)
    shape = add_shape(slide, x, Inches(5.5), card_w, Inches(0.7), BG_CARD, color)
    shape.text_frame.paragraphs[0].text = kw
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].font.color.rgb = color
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


# ================================================================
# SLIDE 2: MCP란 무엇인가?
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "MCP(Model Context Protocol)란?", 36, TEXT_WHITE, True)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
         "AI가 외부 세계와 소통하기 위한 '표준 플러그' - USB처럼 어디든 꽂으면 작동", 18, TEXT_GRAY)

# 비유 설명 카드
add_shape(slide, Inches(0.8), Inches(2), Inches(5.5), Inches(4.5), BG_CARD, ACCENT_BLUE)
add_multiline(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(4), [
    ("쉬운 비유: USB 포트", 22, ACCENT_BLUE, True),
    ("", 8),
    ("과거: 프린터마다 다른 케이블 필요", 16, TEXT_LIGHT),
    ("  → 프린터 A는 케이블 A, 프린터 B는 케이블 B...", 14, TEXT_GRAY),
    ("", 8),
    ("현재: USB 하나로 모든 기기 연결", 16, TEXT_LIGHT),
    ("  → USB 포트 하나면 프린터, 마우스, 키보드 다 OK", 14, TEXT_GRAY),
    ("", 8),
    ("MCP = AI 세계의 USB", 18, ACCENT_GREEN, True),
    ("  → AI가 네이버, 카카오, 공공데이터 등을", 14, TEXT_GRAY),
    ("     하나의 표준 방식으로 연결", 14, TEXT_GRAY),
])

# 기술적 설명 카드
add_shape(slide, Inches(6.8), Inches(2), Inches(5.7), Inches(4.5), BG_CARD, ACCENT_PURPLE)
add_multiline(slide, Inches(7.1), Inches(2.2), Inches(5.2), Inches(4), [
    ("기술적으로는?", 22, ACCENT_PURPLE, True),
    ("", 8),
    ("1. JSON-RPC 2.0 기반 프로토콜", 16, TEXT_LIGHT),
    ("   → AI와 서버가 JSON으로 대화", 14, TEXT_GRAY),
    ("", 6),
    ("2. 3가지 핵심 기능 제공:", 16, TEXT_LIGHT),
    ("   Tools  → AI가 실행할 수 있는 '함수'", 14, ACCENT_BLUE),
    ("   Resources → AI가 읽을 수 있는 '데이터'", 14, ACCENT_GREEN),
    ("   Prompts → 미리 정의된 '대화 템플릿'", 14, ACCENT_ORANGE),
    ("", 6),
    ("3. 전송 방식:", 16, TEXT_LIGHT),
    ("   stdio (로컬) / HTTP+SSE (원격)", 14, TEXT_GRAY),
    ("", 6),
    ("4. Anthropic이 만들고, 오픈소스로 공개", 16, TEXT_LIGHT),
    ("   → Google, OpenAI, MS 모두 채택", 14, ACCENT_GREEN),
])


# ================================================================
# SLIDE 3: MCP 작동 구조
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "MCP는 어떻게 작동하나?", 36, TEXT_WHITE, True)

# 구조도
# Host (Claude/ChatGPT)
add_shape(slide, Inches(0.5), Inches(1.8), Inches(3), Inches(1.5), BG_CARD, ACCENT_BLUE)
add_multiline(slide, Inches(0.7), Inches(1.9), Inches(2.6), Inches(1.3), [
    ("Host (AI 앱)", 20, ACCENT_BLUE, True, PP_ALIGN.CENTER),
    ("Claude Desktop", 14, TEXT_GRAY, False, PP_ALIGN.CENTER),
    ("ChatGPT, Cursor 등", 14, TEXT_GRAY, False, PP_ALIGN.CENTER),
])

# Arrow 1
add_text(slide, Inches(3.5), Inches(2.2), Inches(1.5), Inches(0.8),
         "→ MCP →", 20, ACCENT_GREEN, True, PP_ALIGN.CENTER)

# MCP Client
add_shape(slide, Inches(4.8), Inches(1.8), Inches(3), Inches(1.5), BG_CARD, ACCENT_GREEN)
add_multiline(slide, Inches(5), Inches(1.9), Inches(2.6), Inches(1.3), [
    ("MCP Client", 20, ACCENT_GREEN, True, PP_ALIGN.CENTER),
    ("Host 앱 내부에서", 14, TEXT_GRAY, False, PP_ALIGN.CENTER),
    ("서버와 1:1 연결 관리", 14, TEXT_GRAY, False, PP_ALIGN.CENTER),
])

# Arrow 2
add_text(slide, Inches(7.8), Inches(2.2), Inches(1.5), Inches(0.8),
         "→ JSON →", 20, ACCENT_PURPLE, True, PP_ALIGN.CENTER)

# MCP Server
add_shape(slide, Inches(9.2), Inches(1.8), Inches(3.5), Inches(1.5), BG_CARD, ACCENT_PURPLE)
add_multiline(slide, Inches(9.4), Inches(1.9), Inches(3.1), Inches(1.3), [
    ("MCP Server", 20, ACCENT_PURPLE, True, PP_ALIGN.CENTER),
    ("네이버, 카카오, 공공데이터", 14, TEXT_GRAY, False, PP_ALIGN.CENTER),
    ("등의 API를 감싸는 서버", 14, TEXT_GRAY, False, PP_ALIGN.CENTER),
])

# 아래 설명
add_shape(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(3.2), BG_CARD)

add_multiline(slide, Inches(0.8), Inches(3.9), Inches(11.5), Inches(3), [
    ("실제 작동 예시: '서울 날씨 알려줘'", 22, ACCENT_ORANGE, True),
    ("", 8),
    ("1. 사용자 → Claude에게 '서울 날씨 알려줘' 요청", 16, TEXT_LIGHT),
    ("2. Claude (Host) → MCP Client에게 '날씨 도구 있나?' 확인", 16, TEXT_LIGHT),
    ("3. MCP Client → 기상청 MCP Server에 연결 → get_weather('서울') 호출", 16, TEXT_LIGHT),
    ("4. MCP Server → 기상청 API 호출 → 결과 JSON 반환", 16, TEXT_LIGHT),
    ("5. Claude → 사용자에게 '서울은 현재 맑음, 12도입니다' 답변", 16, TEXT_LIGHT),
    ("", 8),
    ("핵심: MCP Server가 없으면 AI는 기상청 API를 직접 호출할 수 없다!", 18, ACCENT_RED, True),
    ("→ MCP Server = AI와 실제 서비스 사이의 '통역사'", 16, ACCENT_GREEN, True),
])


# ================================================================
# SLIDE 4: 왜 한국 로컬 MCP 서버인가?
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "왜 '한국 로컬' MCP 서버인가?", 36, TEXT_WHITE, True)

# 현황 카드
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5), BG_CARD, ACCENT_RED)
add_multiline(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.2), [
    ("현재 상황 (문제)", 22, ACCENT_RED, True),
    ("", 6),
    ("• 글로벌 MCP 서버: GitHub, Slack, Notion 등 200+개", 15, TEXT_LIGHT),
    ("• 한국 로컬 MCP 서버: 거의 없음!", 15, ACCENT_RED, True),
    ("• 네이버 검색/쇼핑 → MCP 서버 없음", 14, TEXT_GRAY),
    ("• 카카오맵/카카오톡 → MCP 서버 없음", 14, TEXT_GRAY),
    ("• 공공데이터포털 → MCP 서버 없음", 14, TEXT_GRAY),
    ("• 홈택스/국세청 → MCP 서버 없음", 14, TEXT_GRAY),
])

add_shape(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5), BG_CARD, ACCENT_GREEN)
add_multiline(slide, Inches(7.1), Inches(1.7), Inches(5.5), Inches(2.2), [
    ("기회 (우리가 채울 공백)", 22, ACCENT_GREEN, True),
    ("", 6),
    ("• 한국에서 Claude/ChatGPT 사용자 급증", 15, TEXT_LIGHT),
    ("• '클로드야, 네이버에서 맛집 찾아줘' → 불가능", 15, ACCENT_ORANGE),
    ("• 우리가 MCP 서버를 만들면 → 가능해짐", 15, ACCENT_GREEN, True),
    ("• 한국어 + 한국 API 전문가 = 진입장벽", 14, TEXT_LIGHT),
    ("• 영어권 개발자가 따라올 수 없는 영역", 14, TEXT_LIGHT),
])

# 시장 규모
add_shape(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.7), BG_CARD, ACCENT_BLUE)
add_multiline(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.3), [
    ("시장 타이밍이 완벽한 이유", 22, ACCENT_BLUE, True),
    ("", 6),
    ("2025.11  Anthropic, MCP 프로토콜 오픈소스 공개", 15, TEXT_LIGHT),
    ("2025.12  Google, OpenAI도 MCP 지원 발표", 15, TEXT_LIGHT),
    ("2026.01  MCP 서버 수 1,000개 돌파 (대부분 영어권)", 15, TEXT_LIGHT),
    ("2026.03  현재 → 한국 로컬 MCP 서버 시장 '거의 공백'", 15, ACCENT_ORANGE, True),
    ("2026.06  예상 → AI 에이전트 보편화로 MCP 수요 폭증", 15, ACCENT_GREEN),
    ("", 6),
    ("→ 지금 시작하면 6개월 이상 선점 우위 확보 가능", 18, ACCENT_GREEN, True),
])


# ================================================================
# SLIDE 5: 어떤 MCP 서버를 만들 것인가?
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "어떤 MCP 서버를 만들 것인가?", 36, TEXT_WHITE, True)

# Phase 1
add_shape(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.2), BG_CARD, ACCENT_GREEN)
add_multiline(slide, Inches(0.8), Inches(1.7), Inches(3.3), Inches(4.8), [
    ("Phase 1: MVP (1~2주)", 20, ACCENT_GREEN, True),
    ("", 6),
    ("네이버 검색 MCP", 16, TEXT_WHITE, True),
    ("  • 웹 검색, 뉴스, 블로그", 13, TEXT_GRAY),
    ("  • 쇼핑 검색, 이미지 검색", 13, TEXT_GRAY),
    ("", 4),
    ("네이버 지도 MCP", 16, TEXT_WHITE, True),
    ("  • 장소 검색, 길찾기", 13, TEXT_GRAY),
    ("  • 주변 맛집/카페 찾기", 13, TEXT_GRAY),
    ("", 4),
    ("공공데이터 MCP", 16, TEXT_WHITE, True),
    ("  • 날씨 (기상청 API)", 13, TEXT_GRAY),
    ("  • 대중교통 (버스/지하철)", 13, TEXT_GRAY),
    ("  • 사업자등록 확인", 13, TEXT_GRAY),
])

# Phase 2
add_shape(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(5.2), BG_CARD, ACCENT_BLUE)
add_multiline(slide, Inches(4.9), Inches(1.7), Inches(3.3), Inches(4.8), [
    ("Phase 2: 확장 (3~4주)", 20, ACCENT_BLUE, True),
    ("", 6),
    ("카카오 MCP", 16, TEXT_WHITE, True),
    ("  • 카카오맵 장소 검색", 13, TEXT_GRAY),
    ("  • 카카오톡 채널 메시지", 13, TEXT_GRAY),
    ("", 4),
    ("부동산 MCP", 16, TEXT_WHITE, True),
    ("  • 실거래가 조회", 13, TEXT_GRAY),
    ("  • 아파트 시세", 13, TEXT_GRAY),
    ("", 4),
    ("금융 MCP", 16, TEXT_WHITE, True),
    ("  • 주식 시세 (한국거래소)", 13, TEXT_GRAY),
    ("  • 환율 정보", 13, TEXT_GRAY),
    ("  • 은행 금리 비교", 13, TEXT_GRAY),
])

# Phase 3
add_shape(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(5.2), BG_CARD, ACCENT_PURPLE)
add_multiline(slide, Inches(9), Inches(1.7), Inches(3.6), Inches(4.8), [
    ("Phase 3: 프리미엄 (5~8주)", 20, ACCENT_PURPLE, True),
    ("", 6),
    ("홈택스/세무 MCP", 16, TEXT_WHITE, True),
    ("  • 세금 계산기", 13, TEXT_GRAY),
    ("  • 사업자 세무 조회", 13, TEXT_GRAY),
    ("", 4),
    ("쿠팡/쇼핑 MCP", 16, TEXT_WHITE, True),
    ("  • 상품 검색/비교", 13, TEXT_GRAY),
    ("  • 최저가 알림", 13, TEXT_GRAY),
    ("", 4),
    ("허브 웹사이트", 16, TEXT_WHITE, True),
    ("  • MCP 서버 카탈로그", 13, TEXT_GRAY),
    ("  • 원클릭 설치 가이드", 13, TEXT_GRAY),
    ("  • API 키 관리 대시보드", 13, TEXT_GRAY),
    ("  • 사용량 모니터링", 13, TEXT_GRAY),
])


# ================================================================
# SLIDE 6: 수익 모델
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "수익 모델", 36, TEXT_WHITE, True)

# Free tier
add_shape(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(3), BG_CARD, TEXT_GRAY)
add_multiline(slide, Inches(0.8), Inches(1.7), Inches(3.3), Inches(2.6), [
    ("Free (무료)", 22, TEXT_GRAY, True),
    ("$0/월", 28, TEXT_GRAY, True),
    ("", 6),
    ("• 기본 MCP 서버 5개", 14, TEXT_LIGHT),
    ("• 일 100회 API 호출", 14, TEXT_LIGHT),
    ("• 커뮤니티 지원", 14, TEXT_LIGHT),
    ("• 오픈소스 접근", 14, TEXT_LIGHT),
])

# Pro tier
add_shape(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(3), BG_CARD, ACCENT_BLUE)
add_multiline(slide, Inches(4.9), Inches(1.7), Inches(3.3), Inches(2.6), [
    ("Pro (개인)", 22, ACCENT_BLUE, True),
    ("$9.99/월", 28, ACCENT_BLUE, True),
    ("", 6),
    ("• 전체 MCP 서버 접근", 14, TEXT_LIGHT),
    ("• 일 10,000회 API 호출", 14, TEXT_LIGHT),
    ("• 프리미엄 서버 (금융/세무)", 14, TEXT_LIGHT),
    ("• 이메일 지원", 14, TEXT_LIGHT),
])

# Enterprise tier
add_shape(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(3), BG_CARD, ACCENT_PURPLE)
add_multiline(slide, Inches(9), Inches(1.7), Inches(3.6), Inches(2.6), [
    ("Enterprise (기업)", 22, ACCENT_PURPLE, True),
    ("$99/월~", 28, ACCENT_PURPLE, True),
    ("", 6),
    ("• 무제한 API 호출", 14, TEXT_LIGHT),
    ("• 커스텀 MCP 서버 개발", 14, TEXT_LIGHT),
    ("• SLA 보장", 14, TEXT_LIGHT),
    ("• 전담 기술 지원", 14, TEXT_LIGHT),
])

# 수익 전망
add_shape(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2), BG_CARD, ACCENT_GREEN)
add_multiline(slide, Inches(0.8), Inches(5), Inches(11.5), Inches(1.8), [
    ("수익 시뮬레이션 (보수적 추정)", 22, ACCENT_GREEN, True),
    ("", 6),
    ("3개월 후: 무료 500명 + Pro 50명 = 월 $500 (약 70만원)", 15, TEXT_LIGHT),
    ("6개월 후: 무료 2,000명 + Pro 200명 + Enterprise 5개 = 월 $2,500 (약 350만원)", 15, TEXT_LIGHT),
    ("12개월 후: Pro 1,000명 + Enterprise 20개 = 월 $12,000 (약 1,700만원)", 15, ACCENT_GREEN, True),
    ("", 4),
    ("+ 추가 수익: MCP 서버 개발 컨설팅, 기업 맞춤 개발 (건당 500~2,000만원)", 14, ACCENT_ORANGE),
])


# ================================================================
# SLIDE 7: 기술 스택
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "기술 스택 & 아키텍처", 36, TEXT_WHITE, True)

# Backend
add_shape(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.2), BG_CARD, ACCENT_BLUE)
add_multiline(slide, Inches(0.8), Inches(1.7), Inches(3.3), Inches(4.8), [
    ("MCP 서버 (Core)", 20, ACCENT_BLUE, True),
    ("", 6),
    ("언어: TypeScript", 15, TEXT_WHITE, True),
    ("  → Zod 스키마 검증 내장", 13, TEXT_GRAY),
    ("  → 타입 안전성 확보", 13, TEXT_GRAY),
    ("", 4),
    ("SDK: @modelcontextprotocol/sdk", 15, TEXT_WHITE, True),
    ("  → 공식 TypeScript SDK", 13, TEXT_GRAY),
    ("", 4),
    ("전송: stdio + Streamable HTTP", 15, TEXT_WHITE, True),
    ("  → 로컬/원격 모두 지원", 13, TEXT_GRAY),
    ("", 4),
    ("패키지 관리: npm", 15, TEXT_WHITE, True),
    ("  → npx로 원클릭 실행", 13, TEXT_GRAY),
    ("", 4),
    ("배포: npm registry", 15, TEXT_WHITE, True),
    ("  → @korea-mcp/* 네임스페이스", 13, TEXT_GRAY),
])

# Hub Website
add_shape(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(5.2), BG_CARD, ACCENT_GREEN)
add_multiline(slide, Inches(4.9), Inches(1.7), Inches(3.3), Inches(4.8), [
    ("허브 웹사이트", 20, ACCENT_GREEN, True),
    ("", 6),
    ("프레임워크: Next.js 15", 15, TEXT_WHITE, True),
    ("  → App Router", 13, TEXT_GRAY),
    ("  → Server Components", 13, TEXT_GRAY),
    ("", 4),
    ("스타일링: Tailwind CSS", 15, TEXT_WHITE, True),
    ("", 4),
    ("인증: NextAuth.js", 15, TEXT_WHITE, True),
    ("  → GitHub/Google 로그인", 13, TEXT_GRAY),
    ("", 4),
    ("DB: Supabase (PostgreSQL)", 15, TEXT_WHITE, True),
    ("  → 사용량 추적", 13, TEXT_GRAY),
    ("  → API 키 관리", 13, TEXT_GRAY),
    ("", 4),
    ("결제: Stripe / 토스페이먼츠", 15, TEXT_WHITE, True),
])

# DevOps
add_shape(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(5.2), BG_CARD, ACCENT_PURPLE)
add_multiline(slide, Inches(9), Inches(1.7), Inches(3.6), Inches(4.8), [
    ("인프라 & DevOps", 20, ACCENT_PURPLE, True),
    ("", 6),
    ("호스팅: Vercel (웹)", 15, TEXT_WHITE, True),
    ("  → 자동 배포, Edge 함수", 13, TEXT_GRAY),
    ("", 4),
    ("MCP 서버 호스팅:", 15, TEXT_WHITE, True),
    ("  → npm 패키지 배포", 13, TEXT_GRAY),
    ("  → Docker 이미지 (옵션)", 13, TEXT_GRAY),
    ("", 4),
    ("모니터링: Vercel Analytics", 15, TEXT_WHITE, True),
    ("", 4),
    ("CI/CD: GitHub Actions", 15, TEXT_WHITE, True),
    ("  → PR마다 자동 테스트", 13, TEXT_GRAY),
    ("  → 자동 npm 배포", 13, TEXT_GRAY),
    ("", 4),
    ("개발도구: Claude Code", 15, ACCENT_ORANGE, True),
    ("  → 메인 개발 환경", 13, TEXT_GRAY),
])


# ================================================================
# SLIDE 8: 경쟁 분석
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "경쟁 분석 & 차별화", 36, TEXT_WHITE, True)

# 글로벌 경쟁
add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.5), BG_CARD, ACCENT_ORANGE)
add_multiline(slide, Inches(0.8), Inches(1.7), Inches(5.3), Inches(2.2), [
    ("글로벌 MCP 서버 현황", 20, ACCENT_ORANGE, True),
    ("", 6),
    ("• Smithery.ai: MCP 서버 레지스트리 (영어 중심)", 14, TEXT_LIGHT),
    ("• mcp.run: MCP 서버 호스팅 플랫폼", 14, TEXT_LIGHT),
    ("• Composio: 250+ 통합 MCP 서버", 14, TEXT_LIGHT),
    ("• 공식 서버: GitHub, Slack, Google Drive 등 30+", 14, TEXT_LIGHT),
    ("", 4),
    ("→ 전부 영어권 서비스. 한국 로컬 서비스 지원 = 0", 15, ACCENT_RED, True),
])

# 우리의 차별화
add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.5), BG_CARD, ACCENT_GREEN)
add_multiline(slide, Inches(7.1), Inches(1.7), Inches(5.2), Inches(2.2), [
    ("우리의 차별화 포인트", 20, ACCENT_GREEN, True),
    ("", 6),
    ("1. 한국 API 전문성 (네이버/카카오 경험)", 14, TEXT_LIGHT),
    ("2. 한국어 문서 & 가이드 (진입장벽)", 14, TEXT_LIGHT),
    ("3. 한국 결제/세금 시스템 통합", 14, TEXT_LIGHT),
    ("4. 한국 법률/규제 대응 (ISMS 등)", 14, TEXT_LIGHT),
    ("5. 커뮤니티 기반 (한국 개발자 생태계)", 14, TEXT_LIGHT),
    ("", 4),
    ("→ 외국 업체가 따라올 수 없는 로컬 전문성", 15, ACCENT_GREEN, True),
])

# 진입장벽
add_shape(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.7), BG_CARD, ACCENT_BLUE)
add_multiline(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.3), [
    ("우리의 Moat (해자)", 22, ACCENT_BLUE, True),
    ("", 6),
    ("1. 먼저 움직이는 자의 이점 (First Mover)  → 한국 MCP 서버 시장의 표준이 됨", 15, TEXT_LIGHT),
    ("2. 네트워크 효과  → 서버가 많아질수록 허브의 가치 증가", 15, TEXT_LIGHT),
    ("3. 한국어 문서 & 커뮤니티  → 영어를 못하는 한국 개발자들의 유일한 선택지", 15, TEXT_LIGHT),
    ("4. API 키 관리 Lock-in  → 한번 설정하면 이동 비용 발생", 15, TEXT_LIGHT),
    ("", 4),
    ("핵심: 기술이 아니라 '로컬 전문성 + 선점'이 차별화", 16, ACCENT_ORANGE, True),
])


# ================================================================
# SLIDE 9: 실행 로드맵
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "실행 로드맵 (12주 계획)", 36, TEXT_WHITE, True)

# Week 1-2
add_shape(slide, Inches(0.5), Inches(1.5), Inches(2.8), Inches(5), BG_CARD, ACCENT_GREEN)
add_multiline(slide, Inches(0.7), Inches(1.7), Inches(2.5), Inches(4.6), [
    ("Week 1~2", 20, ACCENT_GREEN, True),
    ("MVP 출시", 16, ACCENT_GREEN),
    ("", 6),
    ("□ 프로젝트 셋업", 13, TEXT_LIGHT),
    ("  monorepo 구조", 12, TEXT_GRAY),
    ("□ 네이버 검색 MCP", 13, TEXT_LIGHT),
    ("  웹/뉴스/블로그/쇼핑", 12, TEXT_GRAY),
    ("□ 기상청 날씨 MCP", 13, TEXT_LIGHT),
    ("□ npm 패키지 배포", 13, TEXT_LIGHT),
    ("□ GitHub 공개", 13, TEXT_LIGHT),
    ("□ README 한국어 문서", 13, TEXT_LIGHT),
    ("", 4),
    ("목표: 첫 사용자 확보", 14, ACCENT_GREEN, True),
])

# Week 3-4
add_shape(slide, Inches(3.5), Inches(1.5), Inches(2.8), Inches(5), BG_CARD, ACCENT_BLUE)
add_multiline(slide, Inches(3.7), Inches(1.7), Inches(2.5), Inches(4.6), [
    ("Week 3~4", 20, ACCENT_BLUE, True),
    ("서버 확장", 16, ACCENT_BLUE),
    ("", 6),
    ("□ 카카오맵 MCP", 13, TEXT_LIGHT),
    ("□ 공공데이터 MCP", 13, TEXT_LIGHT),
    ("  교통/버스/지하철", 12, TEXT_GRAY),
    ("□ 네이버 지도 MCP", 13, TEXT_LIGHT),
    ("□ 허브 웹사이트 v1", 13, TEXT_LIGHT),
    ("  서버 목록/검색", 12, TEXT_GRAY),
    ("  설치 가이드", 12, TEXT_GRAY),
    ("□ 블로그 포스트 작성", 13, TEXT_LIGHT),
    ("", 4),
    ("목표: 100명 사용자", 14, ACCENT_BLUE, True),
])

# Week 5-8
add_shape(slide, Inches(6.5), Inches(1.5), Inches(2.8), Inches(5), BG_CARD, ACCENT_PURPLE)
add_multiline(slide, Inches(6.7), Inches(1.7), Inches(2.5), Inches(4.6), [
    ("Week 5~8", 20, ACCENT_PURPLE, True),
    ("수익화", 16, ACCENT_PURPLE),
    ("", 6),
    ("□ 프리미엄 서버 개발", 13, TEXT_LIGHT),
    ("  부동산/금융/세무", 12, TEXT_GRAY),
    ("□ 결제 시스템 연동", 13, TEXT_LIGHT),
    ("□ API 키 관리 시스템", 13, TEXT_LIGHT),
    ("□ 사용량 대시보드", 13, TEXT_LIGHT),
    ("□ Pro 플랜 출시", 13, TEXT_LIGHT),
    ("", 4),
    ("목표: 첫 유료 고객", 14, ACCENT_PURPLE, True),
])

# Week 9-12
add_shape(slide, Inches(9.5), Inches(1.5), Inches(3.3), Inches(5), BG_CARD, ACCENT_ORANGE)
add_multiline(slide, Inches(9.7), Inches(1.7), Inches(2.9), Inches(4.6), [
    ("Week 9~12", 20, ACCENT_ORANGE, True),
    ("성장", 16, ACCENT_ORANGE),
    ("", 6),
    ("□ Enterprise 플랜", 13, TEXT_LIGHT),
    ("□ 커스텀 MCP 개발 서비스", 13, TEXT_LIGHT),
    ("□ 한국 개발자 커뮤니티", 13, TEXT_LIGHT),
    ("  Discord/오픈카톡", 12, TEXT_GRAY),
    ("□ 컨퍼런스 발표", 13, TEXT_LIGHT),
    ("  (if/kakao, DEVIEW 등)", 12, TEXT_GRAY),
    ("□ Smithery.ai 등록", 13, TEXT_LIGHT),
    ("□ 기업 영업 시작", 13, TEXT_LIGHT),
    ("", 4),
    ("목표: MRR $2,000+", 14, ACCENT_ORANGE, True),
])


# ================================================================
# SLIDE 10: 마무리
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(1), Inches(1), Inches(11), Inches(1),
         "핵심 요약", 44, TEXT_WHITE, True, PP_ALIGN.CENTER)

add_shape(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(4), BG_CARD, ACCENT_BLUE)
add_multiline(slide, Inches(2), Inches(2.5), Inches(9), Inches(3.5), [
    ("", 8),
    ('"AI를 사용하는 서비스"가 아닌', 24, TEXT_GRAY, False, PP_ALIGN.CENTER),
    ('"AI가 사용하는 서비스"를 만든다', 28, ACCENT_BLUE, True, PP_ALIGN.CENTER),
    ("", 12),
    ("한국 로컬 API × MCP 프로토콜 × 1인 개발", 20, TEXT_LIGHT, False, PP_ALIGN.CENTER),
    ("= 경쟁자 없는 블루오션", 22, ACCENT_GREEN, True, PP_ALIGN.CENTER),
    ("", 12),
    ("지금 시작하면, 한국 AI 인프라의 표준이 될 수 있다", 20, ACCENT_ORANGE, True, PP_ALIGN.CENTER),
])

# 하단 CTA
add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.6),
         "Let's Build the AI Infrastructure for Korea", 18, TEXT_GRAY, False, PP_ALIGN.CENTER)


# === SAVE ===
output_path = "/home/user/LSY/MCP_서버_허브_사업계획.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
